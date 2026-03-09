from docx import Document
from pptx import Presentation
from pptx.util import Inches

def generate_coach_proposal_docx(coach_data: dict, output_path: str, template_path: str = None):
    """
    Generates a Word document. If template_path is not provided, creates a new one.
    """
    if template_path:
        doc = Document(template_path)
    else:
        doc = Document()
        doc.add_heading('제안서 코치 투입 계획', 0)
        
    doc.add_heading(f"코치: {coach_data.get('name', '이름 없음')}", level=1)
    
    # 텍스트 삽입
    p = doc.add_paragraph()
    p.add_run("전문 분야: ").bold = True
    p.add_run(coach_data.get('domain', 'N/A'))
    
    p = doc.add_paragraph()
    p.add_run("주요 역량: ").bold = True
    p.add_run(coach_data.get('skills', 'N/A'))
    
    # 예산 산출 표 생성 (동적 예산표 모델링)
    # [설계 파이프라인 5단계: 지능형 예산 산출]
    table = doc.add_table(rows=2, cols=4)
    table.style = 'Table Grid'
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = '항목(등급)'
    hdr_cells[1].text = '단가'
    hdr_cells[2].text = '투입수량'
    hdr_cells[3].text = '총액'
    
    row_cells = table.rows[1].cells
    row_cells[0].text = f"메인 코칭비 ({coach_data.get('tier', 'A급')})"
    
    # 임시 단가 시뮬레이션
    base_price = 110000 if coach_data.get('tier') == 'A급' else 80000
    hours = 3
    total = base_price * hours
    
    row_cells[1].text = f"{base_price:,} 원"
    row_cells[2].text = f"{hours} 시간"
    row_cells[3].text = f"{total:,} 원"
    
    doc.add_paragraph(f"산출 근거: {coach_data.get('tier', 'A급')} 전문 코칭비 {base_price:,}원 x {hours}시간 = {total:,}원")
    
    doc.save(output_path)
    return output_path

def generate_coach_profile_pptx(coach_data: dict, image_path: str, output_path: str, template_path: str = None):
    """
    Generates a PowerPoint slide injecting text and the processed grayscale image.
    """
    if template_path:
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
    # 삽입될 슬라이드 (마지막 슬라이드로 가정)
    slide = prs.slides[-1]
    
    # 1. 텍스트 삽입 (이름, 경력)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    tf = txBox.text_frame
    tf.text = f"{coach_data.get('name', '이름 없음')} 코치"
    
    p = tf.add_paragraph()
    p.text = str(coach_data.get('summary', '소개 요약'))
    
    # 2. 전처리된 흑백 사진 삽입 (Placeholder 위치 조정)
    # 실제로는 템플릿의 shape.is_placeholder 를 찾아서 삽입하지만 여기서는 절대좌표 삽입 모의
    if image_path:
        left = Inches(5)
        top = Inches(1)
        # 이미지 크기는 픽셀 해상도 3:4 비율에 맞춰 삽입
        pic = slide.shapes.add_picture(image_path, left, top, width=Inches(3), height=Inches(4))
        
    prs.save(output_path)
    return output_path
